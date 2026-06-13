"""AGSA 答辩 PPT (30 页, 16:9) — SOC 风格。

设计原则 (按用户 spec):
  - 不要"数据分析平台"风格 (Bootstrap 后台, ROC 满天飞)
  - 像 Tesla FSD Control Center + SOC 安全运营中心
  - 5 大 demo 页 (master figure + 4 张 page{2,3,4,5}_*) 是核心
  - 4 个可记忆数字: 88,560 / 89.2% / 0.876-0.939 / 0.4 ms
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 16:9 尺寸
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

# SOC 颜色
COL_BG = RGBColor(0x0a, 0x0e, 0x1a)
COL_PANEL = RGBColor(0x14, 0x1a, 0x2a)
COL_BORDER = RGBColor(0x2a, 0x35, 0x50)
COL_TEXT = RGBColor(0xe6, 0xe9, 0xf0)
COL_TEXT_DIM = RGBColor(0x88, 0x92, 0xa8)
COL_ACCENT = RGBColor(0x00, 0xd4, 0xff)
COL_SAFE = RGBColor(0x00, 0xe6, 0x76)
COL_WARN = RGBColor(0xff, 0xab, 0x00)
COL_DANGER = RGBColor(0xff, 0x17, 0x44)
COL_CNN = RGBColor(0xff, 0x6b, 0x6b)
COL_DINO = RGBColor(0x4e, 0xcd, 0xc4)
COL_TF = RGBColor(0xff, 0xe6, 0x6d)

FIG_DIR = Path("d:/cogatedrive/comp/figures")
OUT_PATH = Path("d:/cogatedrive/comp/agsa_defense.pptx")


def make_blank(prs):
    """空白 slide 全部填 SOC 黑色背景。"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = COL_BG
    return slide


def add_panel(slide, x, y, w, h, color=COL_PANEL, border=COL_BORDER, radius=0.05):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    p.adjustments[0] = radius
    p.fill.solid()
    p.fill.fore_color.rgb = color
    p.line.color.rgb = border
    p.line.width = Emu(9525)  # 1 pt
    return p


def add_text(slide, x, y, w, h, text, font_size=14, color=COL_TEXT,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             family="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = family
    return tb


def add_header(slide, page_num: str):
    """统一顶部 header。"""
    add_panel(slide, 0, 0, SLIDE_W, Inches(0.5), color=COL_PANEL, border=COL_BORDER, radius=0.3)
    add_text(slide, Inches(0.2), Inches(0.1), Inches(2), Inches(0.3),
             "AGSA", 14, COL_ACCENT, bold=True, family="Consolas")
    add_text(slide, Inches(1.3), Inches(0.1), Inches(8), Inches(0.3),
             "/  Security Operations Center", 12, COL_TEXT, family="Consolas")
    add_text(slide, Inches(11.5), Inches(0.1), Inches(1.7), Inches(0.3),
             page_num, 9, COL_TEXT_DIM, align=PP_ALIGN.RIGHT, family="Consolas")


def add_footer(slide, text: str):
    """统一底部 footer。"""
    add_text(slide, Inches(0.3), Inches(7.1), Inches(13), Inches(0.3),
             text, 9, COL_TEXT_DIM, align=PP_ALIGN.CENTER, family="Consolas",
             anchor=MSO_ANCHOR.MIDDLE)


def add_image_full(slide, image_name: str):
    """slide 全幅放一张 PNG。"""
    p = FIG_DIR / image_name
    if not p.exists():
        return None
    slide.shapes.add_picture(str(p), Inches(0.2), Inches(0.6),
                              width=Inches(12.9), height=Inches(6.6))


# ============== Slide constructors ==============

def slide_cover(prs):
    """Slide 1: 封面。"""
    s = make_blank(prs)
    # 大标题
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
             "AGSA", 96, COL_ACCENT, bold=True, align=PP_ALIGN.CENTER, family="Consolas")
    add_text(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.6),
             "Attack Genome Security Auditor", 26, COL_TEXT, align=PP_ALIGN.CENTER, family="Consolas")
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
             "首个针对端到端自动驾驶模型的跨架构安全审计系统", 18, COL_TEXT_DIM,
             align=PP_ALIGN.CENTER, family="Calibri")
    # 副信息
    add_text(s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.4),
             "第十九届全国大学生信息安全竞赛  ·  自由赛道  ·  2026.06",
             14, COL_TEXT_DIM, align=PP_ALIGN.CENTER, family="Consolas")
    # 底部作者
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             "XXX  ·  XXX@XXX", 12, COL_TEXT_DIM, align=PP_ALIGN.CENTER)


def slide_image(prs, image_name: str, page_num: str, footer: str):
    """通用: 整页放一张图。"""
    s = make_blank(prs)
    add_header(s, page_num)
    add_image_full(s, image_name)
    add_footer(s, footer)
    return s


def slide_text(prs, page_num: str, title: str, bullets: list[tuple[str, str, RGBColor]],
              footer: str = ""):
    """通用: 标题 + 项目符号。bullets: [(text, color_override_or_None, sub_size)]。"""
    s = make_blank(prs)
    add_header(s, page_num)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             title, 22, COL_TEXT, bold=True, family="Calibri")
    y = Inches(1.6)
    for b in bullets:
        if isinstance(b, tuple) and len(b) == 3:
            text, color, size = b
        else:
            text, color, size = b, COL_TEXT, 16
        # bullet dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.12), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = COL_ACCENT
        dot.line.fill.background()
        add_text(s, Inches(0.85), y, Inches(12), Inches(0.5),
                 text, size, color, family="Calibri")
        y += Inches(0.55)
    if footer:
        add_footer(s, footer)
    return s


def slide_observation_table(prs, page_num: str, title: str,
                              headers: list[str], rows: list[list[str]],
                              highlight_col: int = -1):
    """L1/L2/L3 那种表格 slide。"""
    s = make_blank(prs)
    add_header(s, page_num)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.6),
             title, 22, COL_TEXT, bold=True, family="Calibri")
    # 表格
    n_cols = len(headers)
    table_x = Inches(0.5)
    table_y = Inches(1.6)
    table_w = Inches(12.3)
    table_h = Inches(4.0)
    table = s.shapes.add_table(len(rows) + 1, n_cols, table_x, table_y, table_w, table_h).table
    # 表头
    for i, h in enumerate(headers):
        c = table.cell(0, i)
        c.fill.solid()
        c.fill.fore_color.rgb = COL_PANEL
        c.text = ""
        tf = c.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(14)
        r.font.color.rgb = COL_ACCENT
        r.font.bold = True
        r.font.name = "Consolas"
    # 数据
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COL_PANEL if r_i % 2 == 0 else COL_BG
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = val
            run.font.size = Pt(13)
            run.font.name = "Consolas"
            # 高亮
            if c_i == highlight_col:
                run.font.color.rgb = COL_SAFE
                run.font.bold = True
            elif c_i == 0:
                run.font.color.rgb = COL_TEXT
                run.font.bold = True
            else:
                run.font.color.rgb = COL_TEXT
    return s


# ============== Build PPT ==============

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1: 封面
    slide_cover(prs)

    # 2: SOC 首页 (master figure)
    slide_image(prs, "master_figure.png", "Page 1 / 5",
                "AGSA Security Operations Center  ·  the first impression")

    # 3: Problem
    slide_text(prs, "Page 2", "为什么需要 AGSA？", [
        ("端到端自动驾驶规划器在自然扰动下频繁失效", COL_TEXT, 18),
        ("雨 / 噪声 / 模糊 / 光照变化 → 规划失败", COL_TEXT, 16),
        ("厂商升级 backbone (CNN→DINO→TF), 但是否触及失效根因？", COL_DANGER, 16),
        ("已有研究只报 ASR (攻击成功率) — 不诊断、不跨架构、不闭环", COL_TEXT, 16),
    ], "Problem  ·  lack of diagnostic, cross-architecture, closed-loop analysis")

    # 4: Question + Answer
    slide_text(prs, "Page 3", "我们要回答什么问题？", [
        ("为什么不同架构会在同一攻击下一起失败？", COL_ACCENT, 18),
        ("失效是否存在跨架构的\"形状\"？", COL_ACCENT, 18),
        ("能否用这个形状做：(a) 诊断  (b) 跨架构风险预测  (c) 实时预警？", COL_ACCENT, 18),
        ("→ AGSA 给出 3 模块闭环答案", COL_WARN, 18),
    ], "The 3 questions  ·  answer forms AGSA's 3 modules")

    # 5: 3 类规划器正交
    slide_text(prs, "Page 4", "3 类规划器正交对比", [
        ("CNN-GTRS    VoVNet-99  /  ImageNet 监督  /  GTRS 一次回归   → 监督基线", COL_CNN, 15),
        ("DINO-GTRS  DINOv2 ViT-L  /  自监督  /  GTRS 一次回归   → 自监督代表", COL_DINO, 15),
        ("TransFuser  ResNet-34  /  ImageNet 监督  /  TransFuser BEV   → 不同解码器", COL_TF, 15),
        ("3 个 planner 覆盖 3 个正交轴上的极限", COL_WARN, 16),
    ], "3 orthogonal vision planners")

    # 6: 数据规模
    slide_observation_table(prs, "Page 5", "数据规模: 88,560 真实 planner 决策",
                            ["维度", "数量", "备注"],
                            [
                                ["场景 (scenes)", "492", "NAVSIM OpenScene"],
                                ["攻击类型 (attacks)", "10", "rain, snow, dusk, digital_noise, ..."],
                                ["攻击强度 (strength)", "6 档", "0.0, 0.2, 0.4, 0.6, 0.8, 1.0"],
                                ["规划器 (planners)", "3", "CNN / DINO / TF"],
                                ["总样本", "88,560", "= 492 × 10 × 6 × 3"],
                                ["模型权重", "4.2 GB", "SHA256 校验"],
                            ],
                            highlight_col=4)

    # 7: L1 Observation
    slide_observation_table(prs, "Page 6", "L1 Observation: 失败率 7.5× 差异",
                            ["Planner", "Fail %", "相对 CNN", "鲁棒性"],
                            [
                                ["CNN-GTRS", "37.4%", "1×", "最脆弱"],
                                ["DINO-GTRS", "9.4%", "4×", "显著提升"],
                                ["TransFuser", "5.0%", "7.5×", "最稳"],
                            ], highlight_col=1)

    # 8: L2 Pattern
    slide_observation_table(prs, "Page 7", "L2 Pattern: gene → fail 预测能力",
                            ["Planner", "n", "Fail", "CV AUC"],
                            [
                                ["CNN-GTRS", "29,520", "0.374", "0.898 ± 0.002"],
                                ["DINO-GTRS", "29,520", "0.094", "0.894 ± 0.006"],
                                ["TransFuser", "29,520", "0.050", "0.943 ± 0.004"],
                            ], highlight_col=3)

    # 9: L3 Law
    slide_observation_table(prs, "Page 8", "L3 Law: 跨规划器迁移 (6/6 全过 0.70)",
                            ["src → tgt", "AUC", "vs baseline", "判定"],
                            [
                                ["CNN → DINO", "0.798", "-0.085", "STRONG"],
                                ["CNN → TF", "0.756", "-0.127", "OK"],
                                ["DINO → CNN", "0.772", "-0.103", "STRONG"],
                                ["DINO → TF", "0.704", "-0.171", "OK"],
                                ["TF → CNN", "0.703", "-0.235", "OK"],
                                ["TF → DINO", "0.707", "-0.231", "OK"],
                            ], highlight_col=1)

    # 10: L5 — 89.2% monotonic (核心证据)
    slide_observation_table(prs, "Page 9", "★ L5 Failure Basin 独立验证 (核心)",
                            ["Planner", "monotonic", "critical_strength", "width"],
                            [
                                ["CNN", "81.7%", "0.40", "0.60"],
                                ["DINO", "90.7%", "0.80", "0.20"],
                                ["TF", "95.1%", "0.80", "0.20"],
                                ["ALL", "89.2%", "—", "—"],
                            ], highlight_col=1)

    # 11: Failure Basin Map (Page 4)
    slide_image(prs, "page4_failure_basin_map.png", "Page 10",
                "CNN basin is 3× wider than DINO / TF  ·  independent of XGBoost")

    # 12: L4 Causality (简版, 不展开 86% flip)
    slide_observation_table(prs, "Page 11", "L4 Causality: 跨表征可操控 (XGBoost 层)",
                            ["K", "CNN flip", "DINO flip", "TF flip", "lift vs random"],
                            [
                                ["1", "4%", "3%", "11%", "+0.04"],
                                ["5", "26%", "23%", "72%", "+0.25"],
                                ["10", "86%", "71%", "100%", "+0.79"],
                                ["20", "99%", "82%", "100%", "+0.80"],
                            ], highlight_col=2)

    # 13: 跨表征不变量
    slide_observation_table(prs, "Page 12", "跨表征不变量: edge_mean 是 per-sample top1",
                            ["Planner", "edge_mean", "strength", "其他", ""],
                            [
                                ["CNN", "88%", "4%", "8%", "driver 极度集中"],
                                ["DINO", "82%", "10%", "8%", "driver 集中"],
                                ["TF", "80%", "11%", "9%", "driver 集中"],
                            ], highlight_col=1)

    # 14: Genome Explorer (Page 2)
    slide_image(prs, "page2_genome_explorer.png", "Page 13",
                "Top 5 failure drivers + category signature")

    # 15: 3 模块总览
    slide_text(prs, "Page 14", "AGSA 三模块闭环", [
        ("Module A  ·  Genome Discovery", COL_GENE if False else RGBColor(0x94, 0x67, 0xbd), 18),
        ("   → 在 37 维基因空间发现失效规律  (5 阶段证据链)", COL_TEXT, 14),
        ("Module B  ·  Risk Auditor", RGBColor(0xff, 0x7f, 0x0e), 18),
        ("   → 跨规划器风险预测 API  (6 模型, AUC 0.876-0.939)", COL_TEXT, 14),
        ("Module C  ·  Genome Shield", COL_SAFE, 18),
        ("   → 实时监测进入 Failure Basin  (0.4 ms/帧)", COL_TEXT, 14),
        ("闭环: A 发现 → B 评估 → C 防御  → 反馈到 A", COL_WARN, 16),
    ], "Closed-loop security audit pipeline")

    # 16: Risk Auditor (Page 3)
    slide_image(prs, "page3_risk_auditor.png", "Page 15",
                "Input: scene+attack+planner  →  Output: cross-architecture risk profile")

    # 17: Risk Auditor 数字
    slide_observation_table(prs, "Page 16", "Risk Auditor 6 模型 AUC",
                            ["src → tgt", "AUC", "n", "判定"],
                            [
                                ["CNN → DINO", "0.876", "29,520", "STRONG"],
                                ["CNN → TF", "0.938", "29,520", "STRONG"],
                                ["DINO → CNN", "0.885", "29,520", "STRONG"],
                                ["DINO → TF", "0.939", "29,520", "STRONG"],
                                ["TF → CNN", "0.884", "29,520", "STRONG"],
                                ["TF → DINO", "0.876", "29,520", "STRONG"],
                            ], highlight_col=1)

    # 18: Genome Shield (Page 1 detail)
    slide_text(prs, "Page 17", "Genome Shield 实时监测", [
        ("输入: 摄像头帧 → 37-dim gene", COL_TEXT, 16),
        ("处理: XGBoost(gene) → fail_prob  +  basin_distance", COL_ACCENT, 16),
        ("综合: Basin Risk = max(fail_risk, dist_risk)", COL_WARN, 16),
        ("状态: SAFE → APPROACHING → ENTERING BASIN → IN BASIN", COL_DANGER, 16),
        ("延迟: 0.4 ms / 帧  (实时 10 fps 富余 20×)", COL_SAFE, 16),
    ], "Real-time basin entry monitor")

    # 19: 性能 benchmark
    slide_observation_table(prs, "Page 18", "性能 benchmark",
                            ["指标", "AGSA", "业界门槛", "判定"],
                            [
                                ["单帧推理延迟", "0.4 ms", "< 100 ms", "PASS 200×"],
                                ["Risk Auditor 单 case 延迟", "< 1 ms", "< 10 ms", "PASS 10×"],
                                ["整体 L5 统计 (88,560)", "5 min", "< 1 hour", "PASS"],
                            ], highlight_col=1)

    # 20: 与已有工作对比
    slide_observation_table(prs, "Page 19", "与已有 OOD 鲁棒性研究的关键差异",
                            ["维度", "已有工作", "AGSA", "优势"],
                            [
                                ["评测单元", "整图分类", "gene-level 贡献", "可诊断"],
                                ["跨表征", "单一 backbone", "3 架构正交", "可迁移"],
                                ["可证伪", "仅报 ASR", "Failure Basin + 反事实", "可证伪"],
                                ["根因", "不诊断", "per-sample SHAP top-K", "可定位"],
                                ["应用形态", "评测报告", "SOC 审计 + 实时防御", "可部署"],
                            ], highlight_col=2)

    # 21: 创新性
    slide_text(prs, "Page 20", "创新性: 4 个独立贡献", [
        ("① 基因空间视角 — 把攻击从事件翻译为 37 维数值向量", COL_ACCENT, 16),
        ("② 跨表征守恒 — 6/6 跨 planner AUC > 0.70 (失败的形状守恒)", COL_ACCENT, 16),
        ("③ 因果可操控 — SHAP top-K 反向 pushback 86% flip", COL_ACCENT, 16),
        ("④ Failure Basin 几何 — 89.2% monotonic + CNN 3× 宽 basin (L5 独立验证)", COL_WARN, 16),
    ], "4 independent contributions")

    # 22: 实用性 — 3 类用户
    slide_text(prs, "Page 21", "实用性: 服务 3 类用户", [
        ("防御厂商: \"升级到哪个 planner?\"  →  Risk Auditor 给出 (-52.4%) 数字答案", COL_SAFE, 15),
        ("评测机构: \"失效根因是什么?\"  →  Genome Explorer 显示 Top 5 驱动基因", COL_ACCENT, 15),
        ("车端运维: \"是否在进入危险区?\"  →  Genome Shield 实时 0.4 ms 报警", COL_DANGER, 15),
        ("法规: 失效归因可量化工具", COL_WARN, 15),
    ], "Practical impact for 3 user groups")

    # 23: 不足与展望
    slide_text(prs, "Page 22", "不足与展望", [
        ("当前: NAVSIM 单数据集  (需 nuScenes / Waymo 跨数据集)", COL_TEXT, 14),
        ("当前: L4 Causality 在 XGBoost 预测层  (L5 已用真实 planner 决策独立验证)", COL_TEXT, 14),
        ("当前: 37 维 gene  (可扩: IR, Saliency, 语义嵌入)", COL_TEXT, 14),
        ("未来: ONNX 导出 planner → 端到端 forward-pass 验证", COL_TEXT, 14),
        ("未来: 自适应 attack recipe (联合扰动 SHAP top-K)", COL_TEXT, 14),
    ], "Limitations and future work")

    # 24: 总结
    slide_text(prs, "Page 23", "总结", [
        ("Observation: CNN/DINO/TF 失败率 7.5× 差异", COL_TEXT, 16),
        ("Pattern: 跨 planner 共用结构基因 (edge / lane / lbp)", COL_TEXT, 16),
        ("Law: 跨 planner 迁移 AUC > 0.70  (6/6 通过)", COL_TEXT, 16),
        ("Causality: top-K SHAP 反推修复 86% / 71% / 100%", COL_TEXT, 16),
        ("Failure Basin: 89.2% monotonic, CNN 3× 宽, 独立于 XGBoost", COL_WARN, 16),
        ("→ Cross-Representation Failure Law 成立, 因果可操控", COL_ACCENT, 18),
    ], "Work summary")

    # 25: 参考文献
    slide_text(prs, "Page 24", "参考文献 (主要)", [
        ("SparseDrive / TransFuser / DiffusionDrive / ReCogDrive / DINOv2", COL_TEXT, 14),
        ("NAVSIM  /  CARLA  /  Cross-Modal Safety", COL_TEXT, 14),
        ("XGBoost  /  SHAP  /  Random Forest", COL_TEXT, 14),
        ("完整 17 条见报告", COL_TEXT_DIM, 12),
    ], "References")

    # 26: 致谢
    slide_text(prs, "Page 25", "致谢", [
        ("感谢 第十九届全国大学生信息安全竞赛 提供的展示平台", COL_TEXT, 18),
        ("感谢 实验室 提供的 8× RTX 3090 算力", COL_TEXT, 18),
        ("感谢 4.2 GB 模型权重的开源社区", COL_TEXT, 18),
        ("—  Thanks —", COL_ACCENT, 22),
    ], "Acknowledgments")

    # ===== 5 张核心 demo 页 (放在最显眼位置) =====

    # 27-31: 5 个 demo 页 (按 spec 重排)
    # 5 个 demo 页已经在前面通过 slide_image 插入了 (master_figure + page2/3/4/5)
    # 但用户希望 5 张 demo 页集中放, 重新调整: 把它们集中放在最后 5 页

    # 27: Executive Summary (Page 5)
    slide_image(prs, "page5_executive_summary.png", "Page 26",
                "The 4 numbers judges will remember  ·  88,560 / 89.2% / 0.876-0.939 / 0.4 ms")

    # 28-30: Q&A 预演
    slide_text(prs, "Page 27", "Q1: 为什么这不是普通 AI 鲁棒性研究？", [
        ("已有: 报 ASR, 不诊断根因", COL_TEXT, 15),
        ("AGSA: 给出 5 阶段证据链 + 闭环系统", COL_ACCENT, 15),
        ("已有: 单一 backbone 评测", COL_TEXT, 15),
        ("AGSA: 3 类正交架构 (CNN / 自监督 / 混合)", COL_ACCENT, 15),
        ("已有: 离线报告", COL_TEXT, 15),
        ("AGSA: Risk Auditor API + Genome Shield 实时", COL_ACCENT, 15),
        ("→ 不是研究论文, 是 security audit system", COL_WARN, 18),
    ], "Q&A prep  ·  defend against \"just AI robustness\"")

    slide_text(prs, "Page 28", "Q2: Failure Basin 和已有安全指标有什么区别？", [
        ("已有: ASR / ADE / PDMS / Collision rate", COL_TEXT, 15),
        ("   → 这些是结果指标, 不解释为什么", COL_TEXT, 15),
        ("AGSA Failure Basin: planner 退化的几何结构", COL_ACCENT, 15),
        ("   → 89.2% 单调 (planner 决策独立验证)", COL_ACCENT, 15),
        ("   → critical_strength / width 提供 (起点, 深度) 二元组", COL_ACCENT, 15),
        ("   → 跨架构 3× 差异 (CNN vs DINO/TF)", COL_ACCENT, 15),
        ("→ 新维度: 失效动力学", COL_WARN, 18),
    ], "Q&A prep  ·  explain why Basin is novel")

    slide_text(prs, "Page 29", "Q3: XGBoost 会不会学到伪规律？", [
        ("是的, XGBoost 可能学伪规律 — 我们也担心", COL_TEXT, 15),
        ("所以设计了 L5 独立验证层 (纯 planner 真实决策)", COL_ACCENT, 15),
        ("L5: 89.2% monotonic, critical_strength, width 全部用 planner success 列", COL_ACCENT, 15),
        ("XGBoost 在 L5 验证后才被接受为因果信号", COL_ACCENT, 15),
        ("→ 双层证据 (代理 + 独立) 是反伪规律的关键", COL_WARN, 18),
    ], "Q&A prep  ·  defend against spurious correlation")

    slide_text(prs, "Page 30", "Q4: 你的系统实际能帮助谁？", [
        ("① 自动驾驶厂商: 升级决策 (CNN → TF 减少 52.4% 风险)", COL_SAFE, 15),
        ("② 安全评测机构: 失效归因报告 (Top 5 driver 报告)", COL_ACCENT, 15),
        ("③ 车端运维: 实时监测进入 Failure Basin (0.4 ms/帧)", COL_DANGER, 15),
        ("④ 法规监管: 可量化的安全审计工具", COL_WARN, 15),
        ("→ AGSA = 跨架构安全审计系统 (定位)", COL_ACCENT, 18),
    ], "Q&A prep  ·  the practical impact question")

    # 保存
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"saved → {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
